from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE
from pptx.dml.color import RGBColor
from pptx.dml.line import LineFormat

name = '柳委員文成'

# TODO: customized size based on number of characters
if len(name) <= 5:
    size = 80
else:
    size = 65

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = Inches(1.5)
top = Inches(1.15)
width = Inches(7)
height = Inches(2.45)

# add text
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = name
txBox.rotation = 180
tf.paragraphs[0].font.size = Pt(60)

txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = name
tf.paragraphs[0].font.size = Pt(60)


# add logo
pic = slide.shapes.add_picture('./logo.png', left, top)
pic.rotation = 180
pic = slide.shapes.add_picture('./logo.png', left+ Inches(0.5), top + Inches(0.5))


shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.5), top + Inches(0.5), width, height
)
shape.fill.background()
shape.line.dash_style = MSO_LINE.LONG_DASH
shape.line.width = Pt(10)

# TODO: use date time as file name
prs.save('test.pptx')